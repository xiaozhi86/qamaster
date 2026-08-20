#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
extract_text.py — 需求文档/图片统一文本抽取（requirement-review 第0阶段降本脚本·最终版）

用途：当输入含图片/扫描件/含图文档时，先把非文本内容转为纯文本并**就地回填**到图片在
文档中的原始位置（保留图文上下文语义），再喂给文本-only 模型评审，避免 400 硬崩。
当前部署为 GLM 文本-only，故主路=本地OCR+文档解析；多模态LLM为可选增强。

就地回填策略（修正"图片文字飘到文末/丢失"缺口）：
  - Word：遍历 body 子元素（保留文档流顺序、含图），图片就地 OCR 占位插入
  - PDF ：文本块 + 图片按 y 坐标交错就地回填；扫描件整页 OCR（天然就位）；
          多栏/复杂版式坐标不可靠 → 退化为"正文留位置标记 + OCR 文末汇总"
  - PPT ：shapes 按坐标排序，Picture 就地 OCR 占位
  - 图片：单图无文档上下文，直接 OCR

规则契约单一事实源：config/input_rules.json（本脚本加载）。
agent 用法：
  python extract_text.py <文件路径>            # 摘要模式（来源/字符数/预处理方式/置信度/图片数/降级标记 + 前 N 字预览）
  python extract_text.py <文件路径> --json     # 结构化 JSON（供 SKILL 第0阶段程序化解析）
  python extract_text.py <文件路径> --full    # 全文（降本模式下不建议）

本脚本是 skill 自带可复用资产，不删除（对齐 case-design scripts 风格）。
"""
import sys
import os
import io
import json

# Windows 控制台默认 cp936，强制 stdout 输出 UTF-8，避免中文乱码
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

RULES_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "config", "input_rules.json")


# ===== 规则加载 =====
def _load_rules():
    if not os.path.exists(RULES_PATH):
        return {}
    try:
        with open(RULES_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


RULES = _load_rules()
ACCEPTED = RULES.get("accepted_formats", {})
MIN_PER_PAGE = RULES.get("min_chars_per_page", 20)
SUMM_CFG = RULES.get("summary", {})
PREVIEW = SUMM_CFG.get("preview_chars", 2000)
OCR_CFG = RULES.get("ocr", {})
OCR_CONF = OCR_CFG.get("confidence_threshold", 0.6)
DEGRADE = RULES.get("degrade", {})
INPLACE = RULES.get("image_inplace", {})
INPLACE_ON = INPLACE.get("enabled", True)
PLACEHOLDER = INPLACE.get("placeholder", "【图片@位置{pos}(置信度{conf})：{text}】")


# ===== OCR 引擎（多引擎逐个尝试：rapidocr → pytesseract+tesseract → None）=====
# rapidocr 2.x 要求 Python<3.13，3.14 环境装不上时自动降级 pytesseract（系统 tesseract 二进制，
# 与 Python 版本无关）。两者都不可用则返回 None，上层降级提示用户给文字说明。
_OCR_ENGINE = None  # ("rapidocr", engine_obj) 或 ("pytesseract", tesseract_cmd) 或 None


def _try_rapidocr():
    try:
        from rapidocr_onnxruntime import RapidOCR
        return ("rapidocr", RapidOCR())
    except Exception:
        return None


def _try_pytesseract():
    """tesseract 是独立系统二进制，与 Python 版本无关。pytesseract 仅作 subprocess 调用。"""
    try:
        import pytesseract
        # 探测 tesseract 可执行文件可用性（import 成功不代表二进制在 PATH）
        import shutil
        if not shutil.which("tesseract") and not getattr(pytesseract, "pytesseract", None):
            # 尝试调一次看是否真的能跑
            try:
                import subprocess
                subprocess.run(["tesseract", "--version"], capture_output=True, check=False, timeout=5)
            except Exception:
                return None
        return ("pytesseract", pytesseract)
    except Exception:
        return None


def _init_ocr_engine():
    """按 engines 顺序逐个尝试，首个可用的胜出。"""
    engines = OCR_CFG.get("engines", ["rapidocr_onnxruntime", "pytesseract"])
    for eng_name in engines:
        if eng_name in ("rapidocr_onnxruntime", "rapidocr"):
            e = _try_rapidocr()
            if e:
                return e
        elif eng_name == "pytesseract":
            e = _try_pytesseract()
            if e:
                return e
    return None


def _ocr_engine():
    global _OCR_ENGINE
    if _OCR_ENGINE is not None:
        return _OCR_ENGINE
    _OCR_ENGINE = _init_ocr_engine()
    return _OCR_ENGINE


def _ocr_image_bytes(blob, ext="png"):
    """对图片二进制做 OCR，返回 (text, confidence, engine_name)。"""
    eng = _ocr_engine()
    if eng is None:
        return "", 0.0, None
    eng_name, eng_obj = eng
    if eng_name == "rapidocr":
        return _ocr_rapidocr(eng_obj, blob)
    if eng_name == "pytesseract":
        return _ocr_pytesseract(eng_obj, blob)
    return "", 0.0, None


def _ocr_rapidocr(engine, blob):
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(blob))
        if img.mode != "RGB":
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        result, _ = engine(buf.getvalue())
    except Exception:
        result, _ = engine(blob)  # 退化：直接传 bytes
    texts, confs = [], []
    for item in (result or []):
        texts.append(item[1])
        if len(item) > 2:
            confs.append(float(item[2]))
    conf = (sum(confs) / len(confs)) if confs else 0.0
    return "\n".join(texts), round(conf, 3), "rapidocr"


def _ocr_pytesseract(pytesseract, blob):
    """用 image_to_data 取每词置信度，过滤 -1（无识别）后平均。"""
    try:
        from PIL import Image
        import csv
        img = Image.open(io.BytesIO(blob))
        lang = OCR_CFG.get("lang", "chi_sim+eng")
        data = pytesseract.image_to_data(img, lang=lang, output_type=pytesseract.Output.DICT)
        texts, confs = [], []
        for i, conf in enumerate(data.get("conf", [])):
            try:
                c = float(conf)
            except (ValueError, TypeError):
                continue
            if c < 0:
                continue  # tesseract 用 -1 标记非文字行
            word = (data.get("text", [""])[i] or "").strip()
            if not word:
                continue
            texts.append(word)
            confs.append(c)
        conf = (sum(confs) / len(confs)) if confs else 0.0
        return " ".join(texts), round(conf, 3), "pytesseract"
    except Exception:
        return "", 0.0, "pytesseract"


def _ocr_pil(img):
    """对 PIL.Image 做 OCR，返回 (text, confidence, engine_name)。"""
    buf = io.BytesIO()
    if img.mode != "RGB":
        img = img.convert("RGB")
    img.save(buf, format="PNG")
    return _ocr_image_bytes(buf.getvalue(), "png")


def _placeholder(pos, conf, text):
    """就地占位符（位置语义保留）。"""
    conf_str = "%.2f" % conf
    return PLACEHOLDER.format(pos=pos, conf=conf_str, text=text or "")


# ===== 类型探测 =====
def detect_kind(path):
    ext = os.path.splitext(path)[1].lower()
    if ext in (".txt", ".md", ".markdown"):
        return "text", ext
    if ext == ".pdf":
        return "pdf", ext
    if ext in (".docx", ".docs"):
        return "word", ext
    if ext in (".pptx",):
        return "ppt", ext
    if ext in (".xlsx", ".xls"):
        return "excel", ext
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
        return "image", ext
    return "unknown", ext


# ===== 文本直读 =====
def extract_text_file(path):
    for enc in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read(), {"方式": "直读", "编码": enc, "图片数": 0}
        except UnicodeDecodeError:
            continue
    return "", {"方式": "直读失败", "降级": DEGRADE.get("ocr_dep_missing", "")}


# ===== Word：遍历 body 子元素，图片就地回填 =====
def _image_blob_from_drawing(drawing_elem, doc):
    """从 w:drawing/a:blip 的 r:embed 关系取图片二进制与扩展名。"""
    try:
        from docx.oxml.ns import qn
    except ImportError:
        return None, None
    for blip in drawing_elem.findall(".//" + qn("a:blip")):
        embed = blip.get(qn("r:embed"))
        if not embed:
            continue
        try:
            rel = doc.part.rels[embed]
            image_part = rel.target_part
            blob = image_part.blob
            ext = os.path.splitext(image_part.partname)[1].lstrip(".").lower() or "png"
            return blob, ext
        except Exception:
            continue
    return None, None


def extract_word_inplace(path):
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return "", {"方式": "Word解析失败", "降级": "缺 python-docx: pip install python-docx"}
    try:
        doc = Document(path)
    except Exception as e:
        return "", {"方式": "Word解析异常", "降级": str(e)}
    parts, img_k, confs = [], 0, []
    # 遍历 body 子元素，保留文档流顺序（含图、表）
    for elem in doc.element.body:
        if elem.tag == qn("w:p"):
            drawings = elem.findall(".//" + qn("w:drawing"))
            text = "".join(t.text or "" for t in elem.findall(".//" + qn("w:t")))
            if drawings and INPLACE_ON:
                # 先出段落文字，再就地占位每张图
                if text.strip():
                    parts.append(text)
                for d in drawings:
                    blob, ext = _image_blob_from_drawing(d, doc)
                    if not blob:
                        continue
                    img_k += 1
                    ocr_text, conf, eng_name = _ocr_image_bytes(blob, ext)
                    if conf:
                        confs.append(conf)
                    parts.append(_placeholder(img_k, conf, ocr_text))
            elif text.strip():
                parts.append(text)
        elif elem.tag == qn("w:tbl"):
            for row in elem.findall(".//" + qn("w:tr")):
                cells = [
                    "".join(t.text or "" for t in c.findall(".//" + qn("w:t")))
                    for c in row.findall(qn("w:tc"))
                ]
                parts.append(" | ".join(cells))
    conf = round(sum(confs) / len(confs), 3) if confs else None
    eng = _ocr_engine()
    eng_name = eng[0] if eng else None
    info = {"方式": "Word就地回填", "图片数": img_k}
    if eng_name and img_k:
        info["OCR引擎"] = eng_name
    if conf is not None:
        info["置信度"] = conf
    if conf is not None and conf < OCR_CONF:
        info["低置信度降级"] = DEGRADE.get("ocr_low_confidence", "").format(threshold=OCR_CONF, engine=eng_name or "未知")
    return "\n".join(parts), info


# ===== PDF：文本块+图片按 y 交错就地回填；扫描件整页 OCR =====
def _is_scan_page(page_text, min_chars):
    """空页或极少字符 → 疑似扫描件/图片型PDF。完全无文字(空)才判扫描件，
    有字但<min_chars仅作疑似标记，避免误判含少量文字的版式页。"""
    t = (page_text or "").strip()
    return len(t) == 0 or len(t) < min_chars


def _image_blob_from_pdf(page, im):
    """从 pdfplumber 图片对象裁剪取二进制（best-effort）。"""
    try:
        crop = page.within_bbox((im["x0"], im["top"], im["x1"], im["bottom"]))
        pil_img = crop.to_image().original
        buf = io.BytesIO()
        pil_img.save(buf, format="PNG")
        return buf.getvalue(), "png"
    except Exception:
        return None, None


def _ocr_pdf_scan_pages(path, scan_pages):
    """扫描件页整页 OCR（天然就位）。"""
    try:
        from pdf2image import convert_from_path
    except ImportError:
        return "", {"降级": DEGRADE.get("pdf_scan_poppler_missing", "")}
    eng = _ocr_engine()
    if eng is None:
        return "", {"降级": DEGRADE.get("ocr_dep_missing", "")}
    try:
        pages = convert_from_path(path, dpi=200)
    except Exception as e:
        return "", {"降级": DEGRADE.get("pdf_scan_poppler_missing", "") + "（%s）" % e}
    out, confs = [], []
    for idx, img in enumerate(pages, 1):
        if idx not in scan_pages:
            continue
        t, conf, _ = _ocr_pil(img)
        if t:
            out.append(t)
        if conf:
            confs.append(conf)
    conf = round(sum(confs) / len(confs), 3) if confs else 0.0
    return "\n".join(out), {"置信度": conf}


def extract_pdf_inplace(path):
    try:
        import pdfplumber
    except ImportError:
        return "", {"方式": "PDF解析失败", "降级": "缺 pdfplumber: pip install pdfplumber"}
    try:
        pages_out, scan_pages, img_k, confs, append_imgs = [], [], 0, [], []
        with pdfplumber.open(path) as pdf:
            for pi, page in enumerate(pdf.pages, 1):
                ptext = page.extract_text() or ""
                if _is_scan_page(ptext, MIN_PER_PAGE):
                    scan_pages.append(pi)
                    continue  # 扫描件页留给下方整页OCR兜底
                # 非扫描页：文本行 + 图片按 y 交错
                events = []
                try:
                    for ln in page.extract_text_lines():
                        events.append(("text", ln.get("top", 0), ln.get("text", "")))
                except Exception:
                    events.append(("text", 0, ptext))
                for im in page.images:
                    events.append(("img", im.get("top", 0), im))
                events.sort(key=lambda e: e[1])
                page_parts = []
                for kind, _, e in events:
                    if kind == "text":
                        if e.strip():
                            page_parts.append(e)
                    else:
                        # 优先就地裁剪OCR；坐标裁剪失败→降级文末汇总+位置标记
                        blob, _ = _image_blob_from_pdf(page, e)
                        if blob is not None:
                            img_k += 1
                            ocr_text, conf, _ = _ocr_image_bytes(blob, "png")
                            if conf:
                                confs.append(conf)
                            page_parts.append(_placeholder(img_k, conf, ocr_text))
                        else:
                            img_k += 1
                            page_parts.append("→见图%d@页%d-y%d（复杂版式，OCR见文末汇总）" % (img_k, pi, int(e.get("top", 0))))
                            append_imgs.append((img_k, pi, blob))
                pages_out.append("\n".join(page_parts))
        # 扫描件页整页OCR（天然就位）
        if scan_pages:
            ocr_text, scan_info = _ocr_pdf_scan_pages(path, scan_pages)
            if ocr_text:
                pages_out.append("【扫描件第%s页·整页OCR：%s】" % (",".join(map(str, scan_pages)), ocr_text))
            elif scan_info.get("降级"):
                pages_out.append("【扫描件第%s页·OCR失败：%s】" % (",".join(map(str, scan_pages)), scan_info["降级"]))
        # 复杂版式降级：文末汇总
        for k, pi, _ in append_imgs:
            pages_out.append("【图%d@页%d：（复杂版式，位置标记见正文）】" % (k, pi))
        conf = round(sum(confs) / len(confs), 3) if confs else None
        eng = _ocr_engine()
        eng_name = eng[0] if eng else None
        info = {"方式": "PDF就地回填", "图片数": img_k, "扫描件页": scan_pages}
        if eng_name and img_k:
            info["OCR引擎"] = eng_name
        if conf is not None:
            info["置信度"] = conf
        if conf is not None and conf < OCR_CONF:
            info["低置信度降级"] = DEGRADE.get("ocr_low_confidence", "").format(threshold=OCR_CONF, engine=eng_name or "未知")
        return "\n".join(pages_out), info
    except Exception as e:
        return "", {"方式": "PDF解析异常", "降级": str(e)}


# ===== PPT：shapes 按坐标排序，Picture 就地回填 =====
def _image_blob_from_pptx(shp):
    try:
        return shp.image.blob, shp.image.ext.lstrip(".").lower()
    except Exception:
        return None, None


def extract_ppt_inplace(path):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return "", {"方式": "PPT解析失败", "降级": "缺 python-pptx: pip install python-pptx"}
    try:
        prs = Presentation(path)
        parts, img_k, confs = [], 0, []
        for si, slide in enumerate(prs.slides, 1):
            # 按 (top, left) 排序保留视觉顺序
            shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
            for shp in shapes:
                if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and INPLACE_ON:
                    img_k += 1
                    blob_ext = _image_blob_from_pptx(shp)
                    if not blob_ext or not blob_ext[0]:
                        continue
                    blob, ext = blob_ext
                    t, conf, _ = _ocr_image_bytes(blob, ext)
                    if conf:
                        confs.append(conf)
                    parts.append(_placeholder(img_k, conf, t))
                elif shp.has_text_frame and shp.text_frame.text.strip():
                    parts.append(shp.text_frame.text)
        conf = round(sum(confs) / len(confs), 3) if confs else None
        eng = _ocr_engine()
        eng_name = eng[0] if eng else None
        info = {"方式": "PPT就地回填", "图片数": img_k}
        if eng_name and img_k:
            info["OCR引擎"] = eng_name
        if conf is not None:
            info["置信度"] = conf
        if conf is not None and conf < OCR_CONF:
            info["低置信度降级"] = DEGRADE.get("ocr_low_confidence", "").format(threshold=OCR_CONF, engine=eng_name or "未知")
        return "\n".join(parts), info
    except Exception as e:
        return "", {"方式": "PPT解析异常", "降级": str(e)}


# ===== Excel：openpyxl 抽单元格（通常无内嵌图）=====
def extract_excel(path):
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "", {"方式": "Excel解析失败", "降级": "缺 openpyxl: pip install openpyxl"}
    try:
        wb = load_workbook(path, data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts), {"方式": "Excel抽取", "图片数": 0}
    except Exception as e:
        return "", {"方式": "Excel解析异常", "降级": str(e)}


# ===== 单图：直接 OCR =====
def ocr_single_image(path):
    with open(path, "rb") as f:
        blob = f.read()
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    text, conf, eng_name = _ocr_image_bytes(blob, ext)
    info = {"方式": "OCR(%s)" % (eng_name or "无引擎"), "置信度": conf, "图片数": 1}
    if eng_name:
        info["OCR引擎"] = eng_name
    if conf and conf < OCR_CONF:
        info["低置信度降级"] = DEGRADE.get("ocr_low_confidence", "").format(threshold=OCR_CONF, engine=eng_name or "未知")
    return text, info


# ===== 统一入口 =====
def extract(path):
    if not os.path.exists(path):
        return "", {"方式": "失败", "降级": "文件不存在: %s" % path}
    kind, ext = detect_kind(path)
    # 图片类但OCR引擎不可用 → 提前降级（与 ocr_single_image 的 info 字段一致）
    if kind == "image" and _ocr_engine() is None:
        return "", {"方式": "OCR失败", "图片数": 1, "置信度": 0.0, "降级": DEGRADE.get("ocr_dep_missing", "")}
    if kind == "text":
        return extract_text_file(path)
    if kind == "pdf":
        return extract_pdf_inplace(path)
    if kind == "word":
        return extract_word_inplace(path)
    if kind == "ppt":
        return extract_ppt_inplace(path)
    if kind == "excel":
        return extract_excel(path)
    if kind == "image":
        return ocr_single_image(path)
    return "", {"方式": "失败", "降级": "不支持的格式: %s" % ext}


def main():
    if len(sys.argv) < 2:
        print("用法: python extract_text.py <文件路径> [--json|--full]")
        print("  摘要模式(默认)：来源/字符数/预处理方式/置信度/图片数/降级标记 + 前 N 字预览")
        print("  --json  ：结构化 JSON（供 SKILL 第0阶段程序化解析）")
        print("  --full  ：全文（降本模式下不建议）")
        return 1
    path = sys.argv[1]
    mode = "preview"
    if "--json" in sys.argv:
        mode = "json"
    elif "--full" in sys.argv:
        mode = "full"
    text, info = extract(path)
    if mode == "json":
        payload = {
            "来源": os.path.basename(path),
            "字符数": len(text),
            "预处理方式": info.get("方式", ""),
            "置信度": info.get("置信度"),
            "图片数": info.get("图片数", 0),
            "降级标记": info.get("降级") or info.get("低置信度降级", ""),
            "摘要": text[:PREVIEW],
        }
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0
    if mode == "full":
        print(text)
        return 0
    print("===== 文本抽取摘要 =====")
    print("来源: %s" % os.path.basename(path))
    print("字符数: %d" % len(text))
    for k, v in info.items():
        print("%s: %s" % (k, v))
    print("---- 摘要(前%d字) ----" % PREVIEW)
    print(text[:PREVIEW])
    print("========================")
    return 0


if __name__ == "__main__":
    sys.exit(main())
